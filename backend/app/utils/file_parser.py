import fitz  # PyMuPDF
import docx
from pptx import Presentation
import openpyxl
import io

def parse_file_to_text(file_bytes: bytes, filename: str) -> str:
    """
    解析不同类型的文件为纯文本
    """
    ext = filename.split('.')[-1].lower()
    
    if ext in ['txt', 'md', 'csv']:
        return file_bytes.decode('utf-8', errors='ignore')
        
    elif ext == 'pdf':
        text = ""
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            text += page.get_text()
        return text
        
    elif ext == 'docx':
        text = ""
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
        
    elif ext == 'pptx':
        text = ""
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
        
    elif ext == 'xlsx':
        text = ""
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        for sheet_name in wb.sheetnames:
            text += f"--- Sheet: {sheet_name} ---\n"
            sheet = wb[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                text += row_text + "\n"
        return text
        
    else:
        raise ValueError(f"Unsupported file extension: {ext}")
